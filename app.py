import json
import os
import re
from datetime import datetime
from functools import wraps
from pathlib import Path

try:
    import google.generativeai as genai
except Exception:
    genai = None

try:
    from openai import OpenAI
except Exception:
    OpenAI = None

try:
    import requests
    from bs4 import BeautifulSoup
except Exception:
    requests = None
    BeautifulSoup = None

try:
    from PyPDF2 import PdfReader
except Exception:
    PdfReader = None

try:
    from docx import Document
except Exception:
    Document = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

from dotenv import load_dotenv
from flask import Flask, flash, jsonify, redirect, render_template, request, session, url_for, send_from_directory
from werkzeug.security import check_password_hash, generate_password_hash
from werkzeug.utils import secure_filename

from database import get_db, init_app

load_dotenv()

BASE_DIR = Path(__file__).resolve().parent

# Use /tmp on Vercel because the project directory is read-only
if os.getenv("VERCEL"):
    UPLOAD_DIR = Path("/tmp/uploads")
else:
    UPLOAD_DIR = BASE_DIR / "static" / "uploads"

UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

ALLOWED_EXTENSIONS = {"pdf", "doc", "docx", "ppt", "pptx", "txt", "mp4", "webm", "mov", "png", "jpg", "jpeg"}
TEXT_EXTENSIONS = {"pdf", "docx", "pptx", "txt"}


def allowed_file(filename: str) -> bool:
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS


def safe_limit(text: str, limit: int = 12000) -> str:
    text = " ".join((text or "").split())
    return text[:limit]


def create_app():
    app = Flask(__name__)
    app.config["SECRET_KEY"] = os.getenv("SECRET_KEY", "accesslearn-dev-secret-change-before-deploy")
    app.config["UPLOAD_FOLDER"] = str(UPLOAD_DIR)
    app.config["MAX_CONTENT_LENGTH"] = 60 * 1024 * 1024
    init_app(app)

    def ensure_runtime_schema():
        db = get_db()
        with open(BASE_DIR / "schema.sql", "r", encoding="utf-8") as f:
            db.executescript(f.read())
        migrations = [
            "ALTER TABLE notes ADD COLUMN file_filename TEXT",
            "ALTER TABLE notes ADD COLUMN file_original_name TEXT",
            "ALTER TABLE notes ADD COLUMN file_type TEXT",
            "ALTER TABLE notes ADD COLUMN resource_link TEXT",
            "ALTER TABLE notes ADD COLUMN video_link TEXT",
            "ALTER TABLE chat_logs ADD COLUMN note_id INTEGER",
            "ALTER TABLE summaries ADD COLUMN source_type TEXT",
            "ALTER TABLE summaries ADD COLUMN source_title TEXT",
        ]
        for sql in migrations:
            try:
                db.execute(sql)
            except Exception:
                pass
        db.commit()

    with app.app_context():
        ensure_runtime_schema()

    def current_user():
        user_id = session.get("user_id")
        if not user_id:
            return None
        return get_db().execute("SELECT * FROM users WHERE id = ?", (user_id,)).fetchone()

    @app.context_processor
    def inject_user():
        return {"current_user": current_user()}

    def login_required(view):
        @wraps(view)
        def wrapped(*args, **kwargs):
            if not session.get("user_id"):
                flash("Please login to continue.", "warning")
                return redirect(url_for("login"))
            return view(*args, **kwargs)
        return wrapped

    def role_required(role):
        def decorator(view):
            @wraps(view)
            def wrapped(*args, **kwargs):
                user = current_user()
                if not user:
                    return redirect(url_for("login"))
                if user["role"] != role:
                    flash("You do not have permission to access this page.", "danger")
                    return redirect(url_for("dashboard"))
                return view(*args, **kwargs)
            return wrapped
        return decorator

    # ---------- text extraction ----------
    def extract_file_text(filename: str) -> str:
        if not filename:
            return ""
        path = UPLOAD_DIR / filename
        if not path.exists():
            return ""
        ext = path.suffix.lower().replace('.', '')
        try:
            if ext == "txt":
                return path.read_text(encoding="utf-8", errors="ignore")
            if ext == "pdf" and PdfReader:
                reader = PdfReader(str(path))
                return "\n".join(page.extract_text() or "" for page in reader.pages[:20])
            if ext == "docx" and Document:
                doc = Document(str(path))
                return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            if ext == "pptx" and Presentation:
                prs = Presentation(str(path))
                texts = []
                for slide in prs.slides[:30]:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            texts.append(shape.text)
                return "\n".join(texts)
        except Exception as e:
            print("FILE TEXT EXTRACTION ERROR:", e)
            return ""
        return ""

    def fetch_link_text(url: str) -> str:
        if not url or requests is None:
            return ""
        url = (url or "").strip()
        if not url or "example.com" in url:
            return ""
        try:
            headers = {
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 Chrome/120 Safari/537.36",
                "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
            }
            res = requests.get(url, headers=headers, timeout=15, allow_redirects=True)
            if not res.ok:
                return ""
            content_type = res.headers.get("content-type", "").lower()

            # Normal webpages
            if "html" in content_type or "text" in content_type or not content_type:
                if BeautifulSoup:
                    soup = BeautifulSoup(res.text, "html.parser")
                    for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
                        tag.decompose()
                    text = soup.get_text(" ", strip=True)
                    return safe_limit(text, 20000)
                return safe_limit(res.text, 20000)

            return ""
        except Exception as e:
            print("LINK READ ERROR:", e)
            return ""

    def note_context(note_id: int) -> tuple[str, str]:
        if not note_id:
            return "", ""
        note = get_db().execute("SELECT * FROM notes WHERE id=?", (note_id,)).fetchone()
        if not note:
            return "", ""
        pieces = [note["title"] or "", note["subject"] or "", note["description"] or "", note["content"] or ""]
        file_text = extract_file_text(note["file_filename"] or "")
        if file_text:
            pieces.append(file_text)
        text = safe_limit("\n".join(pieces), 10000)
        return text, note["title"]

    # ---------- AI + fallback engines ----------
    EDUCATIONAL_TERMS = {
        "definition", "formula", "concept", "example", "application", "advantages", "disadvantages",
        "types", "process", "function", "algorithm", "database", "programming", "computer",
        "machine learning", "deep learning", "artificial intelligence", "voltage", "current",
        "resistance", "normalization", "primary key", "foreign key", "syntax", "variable",
        "data structure", "array", "stack", "queue", "network", "operating system", "security"
    }
    PROMO_TERMS = {
        "bhaiya", "didi", "hype", "dominate", "not your college", "nyc", "brand", "marketing",
        "subscribe", "pricing", "sales", "ambitious students", "tier-2", "tier-3", "not in your syllabus"
    }

    def clean_lines(text: str) -> str:
        return re.sub(r"\s+", " ", text or "").strip()

    def is_educational_content(text: str) -> bool:
        lower = (text or "").lower()
        if len(clean_lines(text)) < 120:
            return False
        promo_hits = sum(1 for t in PROMO_TERMS if t in lower)
        edu_hits = sum(1 for t in EDUCATIONAL_TERMS if t in lower)
        # Promotional notes should not be treated as study notes even if they mention AI/coding.
        if promo_hits >= 2 and edu_hits < 5:
            return False
        return edu_hits >= 2 or len(re.findall(r"[.!?]", text)) >= 4

    def detect_topic(text: str) -> str:
        lower = (text or "").lower()
        if any(k in lower for k in ["ohm", "voltage", "current", "resistance"]):
            return "ohm"
        if any(k in lower for k in ["dbms", "database", "sql", "normalization", "primary key"]):
            return "dbms"
        if any(k in lower for k in ["python", "list", "tuple", "dictionary", "programming"]):
            return "python"
        if any(k in lower for k in ["artificial intelligence", "machine learning", "deep learning", "neural", " ai "]):
            return "ai"
        if any(k in lower for k in ["data structure", "array", "stack", "queue", "linked list", "tree"]):
            return "ds"
        if any(k in lower for k in ["operating system", "process", "thread", "deadlock", "memory management"]):
            return "os"
        if any(k in lower for k in ["network", "tcp", "ip", "osi", "router"]):
            return "network"
        return "general"

    def important_sentences(text: str, limit: int = 5):
        text = safe_limit(text, 9000)
        sentences = [x.strip() for x in re.split(r"(?<=[.!?])\s+|\n+", text) if len(x.strip()) > 25]
        scored = []
        keywords = ["definition", "important", "type", "application", "advantage", "example", "process", "function", "feature", "formula", "uses", "means", "called", "allows", "helps"]
        for sent in sentences:
            low = sent.lower()
            if any(p in low for p in PROMO_TERMS):
                continue
            score = sum(1 for k in keywords if k in low) + min(len(sent) / 140, 2)
            scored.append((score, sent))
        scored.sort(reverse=True, key=lambda x: x[0])
        picked = []
        for _, sent in scored:
            if sent not in picked:
                picked.append(sent)
            if len(picked) >= limit:
                break
        return picked or sentences[:limit]

    def extract_keywords(text: str, limit: int = 8):
        stop = {"about", "after", "again", "because", "before", "between", "could", "every", "first", "from", "have", "into", "learn", "learning", "more", "notes", "other", "should", "students", "system", "their", "there", "these", "this", "through", "using", "which", "while", "with", "would", "college", "generation", "ambitious", "bhaiya", "didi", "ready", "focus", "actually"}
        words = re.findall(r"[A-Za-z][A-Za-z\-]{3,}", text or "")
        freq = {}
        for w in words:
            k = w.lower().strip("-")
            if k not in stop and k not in PROMO_TERMS:
                freq[k] = freq.get(k, 0) + 1
        ranked = sorted(freq.items(), key=lambda x: (-x[1], x[0]))
        return [w.title() for w, _ in ranked[:limit]]

    def topic_answer(topic: str, message: str = "") -> str:
        t = detect_topic((topic or "") + " " + (message or ""))
        if t == "ohm":
            return """Ohm's Law explains the relationship between voltage, current, and resistance in an electrical circuit.

Formula:
V = I × R

Where:
• V = Voltage in volts
• I = Current in amperes
• R = Resistance in ohms

Example:
If current is 2A and resistance is 5Ω, voltage = 2 × 5 = 10V.

Why it matters:
It helps students calculate circuit values and understand basic electronics."""
        if t == "dbms":
            return """DBMS means Database Management System. It is software used to store, organize, secure, and retrieve data efficiently.

Key concepts:
• Tables store data in rows and columns
• Primary key uniquely identifies a record
• SQL is used to query data
• Normalization reduces duplication
• Transactions maintain reliability

Examples:
SQLite, MySQL, PostgreSQL, Oracle.

Use in AccessLearn:
The database stores users, notes, quiz scores, summaries, and study reports."""
        if t == "python":
            return """Python is a high-level programming language known for simple syntax and readability.

Key concepts:
• Variables store values
• Lists and dictionaries store collections
• Functions organize reusable code
• Loops repeat tasks
• Libraries extend Python for AI, web, and data science

Example:
print('Hello AccessLearn')

Applications:
Web development, automation, AI/ML, data analysis, and scripting."""
        if t == "ai":
            return """Artificial Intelligence (AI) is a branch of computer science that enables machines to perform tasks that normally require human intelligence.

Key concepts:
• Machine Learning: systems learn from data
• Deep Learning: neural networks solve complex problems
• NLP: computers understand human language
• Computer Vision: systems understand images/videos

Examples:
Chatbots, recommendation systems, voice assistants, healthcare diagnosis, and self-driving vehicles.

In education:
AI helps with doubt solving, summaries, quizzes, and personalized learning."""
        if t == "ds":
            return """Data Structures are ways of organizing data so that it can be used efficiently.

Important types:
• Array: stores elements in continuous memory
• Stack: follows LIFO, last in first out
• Queue: follows FIFO, first in first out
• Linked List: nodes connected by links
• Tree/Graph: used for hierarchical and network data

Why important:
They improve searching, sorting, memory usage, and algorithm performance."""
        if t == "os":
            return """An Operating System is system software that manages computer hardware and software resources.

Key concepts:
• Process management
• Memory management
• File management
• Device management
• Security and user interface

Examples:
Windows, Linux, macOS, Android.

It acts as a bridge between users, applications, and hardware."""
        if t == "network":
            return """Computer Networking connects computers and devices so they can share data and resources.

Key concepts:
• IP address identifies a device
• TCP ensures reliable delivery
• Router forwards data between networks
• OSI model explains communication layers
• DNS converts domain names to IP addresses

Applications:
Internet, email, cloud services, video calls, and online learning platforms."""
        return f"""{message or topic}

Simple explanation:
This topic should be understood by learning its definition, key concepts, examples, applications, and practice questions.

Study approach:
• Start with the basic meaning.
• Identify 4–5 important terms.
• Study one real-life example.
• Revise the concept in your own words.
• Attempt a short quiz to check understanding."""

    def note_relevance_score(question: str, context: str) -> int:
        qwords = {w.lower() for w in re.findall(r"[A-Za-z][A-Za-z]{3,}", question or "")}
        c = (context or "").lower()
        return sum(1 for w in qwords if w in c)

    def fallback_answer(message: str, context: str = "") -> str:
        if context:
            if not is_educational_content(context):
                return """Selected note check:
The selected teacher note does not contain enough proper educational content for a notes-based answer.

I will answer generally instead:

""" + topic_answer(message, message)
            if note_relevance_score(message, context) == 0 and len(message.split()) > 2:
                return """Selected note check:
The selected note does not look directly related to your question, so here is a general learning answer:

""" + topic_answer(message, message)
            points = important_sentences(context, 5)
            keywords = extract_keywords(context, 7)
            answer_intro = topic_answer(context, message).split("\n\n")[0]
            return """Notes-Based Learning Answer

Question: {message}

Short Answer:
{answer_intro}

From Selected Notes:
{points}

Important Terms:
{terms}

What to do next:
• Read the above points once.
• Ask one follow-up doubt in the chatbot.
• Generate a quiz from this note for practice.""".format(
                message=message,
                answer_intro=answer_intro,
                points="\n".join(f"• {p}" for p in points),
                terms="\n".join(f"• {t}" for t in (keywords or ["Definition", "Examples", "Applications"])),
            )
        return topic_answer(message, message)

    def call_openai(prompt: str) -> str:
        """Primary AI provider: OpenAI / ChatGPT.
        Put OPENAI_API_KEY in .env or hosting variables.
        """
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key or OpenAI is None:
            return ""
        try:
            client = OpenAI(api_key=api_key)
            model = os.getenv("OPENAI_MODEL", "gpt-4o-mini")
            response = client.chat.completions.create(
                model=model,
                messages=[
                    {
                        "role": "system",
                        "content": (
                            "You are AccessLearn AI, a helpful educational assistant for students and teachers. "
                            "Give accurate, clear, structured, student-friendly answers. "
                            "Never give generic study-method templates unless the user asks for study methods."
                        ),
                    },
                    {"role": "user", "content": prompt},
                ],
                temperature=0.35,
            )
            return response.choices[0].message.content.strip()
        except Exception:
            return ""

    def call_gemini(prompt: str) -> str:
        """Backup AI provider: Gemini."""
        api_key = os.getenv("GEMINI_API_KEY")
        if not api_key or genai is None:
            return ""
        try:
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel(os.getenv("GEMINI_MODEL", "gemini-2.0-flash"))
            response = model.generate_content(prompt)
            return getattr(response, "text", "") or ""
        except Exception:
            return ""

    def call_ai(prompt: str) -> str:
        """AI priority: OpenAI first, Gemini second, local fallback last."""
        openai_text = call_openai(prompt)
        if openai_text:
            print("AccessLearn AI: OpenAI response used")
            return openai_text
        gemini_text = call_gemini(prompt)
        if gemini_text:
            print("AccessLearn AI: Gemini response used")
            return gemini_text
        print("AccessLearn AI: Local fallback used")
        return ""

    def local_explain(question: str, context: str = "") -> str:
        """Last-resort fallback with topic-specific answers; avoids same generic output."""
        q = (question or "").lower()
        combined = (question + " " + (context or "")).lower()
        topic = detect_topic(combined)

        if "ram" in q and "rom" in q:
            return """RAM and ROM Difference

RAM (Random Access Memory):
• Temporary memory used while the computer is running.
• Data is lost when power is turned off.
• It is faster and used for active programs.

ROM (Read Only Memory):
• Permanent memory that stores essential startup instructions.
• Data is not lost when power is turned off.
• It is used to boot and initialize the system.

In short: RAM is temporary working memory, while ROM is permanent startup memory."""

        if topic == "ohm":
            return """Ohm's Law

Ohm's Law states that the voltage across a conductor is directly proportional to the current flowing through it, if temperature and physical conditions remain constant.

Formula: V = I × R

Where:
• V = Voltage in volts
• I = Current in amperes
• R = Resistance in ohms

Example:
If current is 2 A and resistance is 5 Ω, then V = 2 × 5 = 10 V.

Use: Ohm's Law is used to calculate voltage, current, or resistance in electrical circuits."""

        if topic == "ai":
            return """Artificial Intelligence (AI)

Artificial Intelligence is a branch of computer science that enables machines to perform tasks that normally require human intelligence.

Key ideas:
• Learning from data
• Reasoning and decision-making
• Understanding language
• Recognizing images and patterns

Examples:
• Chatbots
• Recommendation systems
• Self-driving cars
• Healthcare diagnosis
• AI tutors in education

AI is useful in education because it can support doubt-solving, personalized learning, summarization, and quiz generation."""

        if topic == "dbms":
            return """DBMS (Database Management System)

A DBMS is software used to store, organize, manage, and retrieve data efficiently.

Key concepts:
• Table: stores data in rows and columns
• Primary key: uniquely identifies each record
• Foreign key: connects two tables
• SQL: language used to query databases
• Normalization: reduces duplicate data

Examples: SQLite, MySQL, PostgreSQL, Oracle.

DBMS is important because it improves data security, consistency, and fast retrieval."""

        if topic == "python":
            return """Python Programming

Python is a high-level programming language known for simple syntax and readability.

Important features:
• Easy to learn
• Supports object-oriented programming
• Large library support
• Used in AI, web development, automation, and data science

Example:
name = "AccessLearn"
print(name)

Python is useful for beginners because its syntax is close to English and it supports many real-world applications."""

        if context and is_educational_content(context):
            points = important_sentences(context, 5)
            terms = extract_keywords(context, 6)
            return """Notes-Based Answer

Question: {q}

Based on the selected teacher notes, the important points are:
{points}

Important terms:
{terms}

Detailed explanation:
The selected notes mainly discuss these concepts. Read the key points first, then revise the important terms and attempt a short quiz to check understanding.""".format(
                q=question,
                points="\n".join(f"• {p}" for p in points) if points else "• The selected notes need more educational detail.",
                terms="\n".join(f"• {t}" for t in terms) if terms else "• Definition\n• Examples\n• Applications",
            )

        return f"""I can answer this better if you ask a specific topic question.

Your question: {question}

Try asking like:
• What is Artificial Intelligence?
• Explain DBMS with examples.
• Difference between RAM and ROM.
• Generate exam notes on Python.
• Explain this selected note in simple words."""

    def ai_answer(message, context=""):
        message = (message or "").strip()
        if not message:
            return "Please ask a clear question."

        usable_context = context if is_educational_content(context) else ""
        context_instruction = ""
        if context and not usable_context:
            context_instruction = "The selected note appears promotional, incomplete, or not educational. Do not use it as factual study context. Answer generally if possible."
        elif usable_context:
            context_instruction = "Use the teacher notes as context, but do not copy them blindly. Explain clearly and add correct educational detail where needed."

        prompt = f"""
You are AccessLearn AI Assistant.

Task:
Answer the student's question accurately and in detail.

Rules:
1. First understand what the user is asking.
2. If it is a general question, answer generally.
3. If it is an educational question, explain concepts with examples.
4. If teacher notes are provided and relevant, use them as context.
5. If notes are irrelevant or promotional, ignore them and answer from knowledge.
6. Do NOT give the same generic study-method response for every question.
7. NEVER write one long paragraph.
8. Always use separate lines, headings, numbering, and bullet points.
9. Use this structure whenever possible:

# Topic Name

## 1. Definition
Write a clear 2-3 line definition.

## 2. Key Concepts
• Point 1
• Point 2
• Point 3

## 3. Important Points
• Point 1
• Point 2
• Point 3

## 4. Example
Give one practical student-friendly example.

## 5. Applications
• Use 1
• Use 2
• Use 3

## 6. Quick Revision
• Short exam-friendly points.

{context_instruction}

Teacher Notes Context:
{safe_limit(usable_context, 9000)}

Student Question:
{message}
"""
        ai_text = call_ai(prompt)
        return ai_text or local_explain(message, usable_context)

    def summarize_text(text):
        clean_text = safe_limit(text, 18000)
        if not clean_text:
            return "Please enter notes to summarize."

        if not is_educational_content(clean_text):
            return """Content Quality Check

This content does not look like proper educational notes. It appears promotional, incomplete, or unrelated to study material.

Please upload/paste proper content such as:
• Teacher notes
• Textbook explanation
• Lecture content
• Topic-based PDF/DOC notes
• Definitions, concepts, examples, and applications

AccessLearn avoids creating misleading summaries from low-quality content."""

        prompt = f"""
You are AccessLearn Smart Notes Summarizer.

Create a useful study summary from the notes below. Do not copy-paste large text. Understand the content and rewrite it for students.

Required output format:

Topic Overview:
• Explain what the topic is in simple words.

Key Concepts:
• List the main concepts from the notes.

Important Definitions:
• Give short definitions of important terms.

Applications / Examples:
• Give examples or real-life applications.

Exam Revision Notes:
• Give points useful for exams.

Practice Questions:
Generate 10 meaningful questions from the content.

Final 3-Line Summary:
• Summarize the whole topic in 3 lines.

Notes:
{clean_text}
"""
        ai_summary = call_ai(prompt)
        if ai_summary:
            return ai_summary

        points = important_sentences(clean_text, 7)
        terms = extract_keywords(clean_text, 8)
        topic = detect_topic(clean_text)
        overview = local_explain(topic if topic != "general" else (terms[0] if terms else "this topic"), clean_text).split("\n\n")[0]
        return """Topic Overview:
• {overview}

Key Concepts:
{points}

Important Definitions:
{terms}

Applications / Examples:
• This topic can be applied in academic learning, practical projects, exam preparation, and real-world problem solving.

Exam Revision Notes:
• Learn the definition first.
• Revise key concepts and examples.
• Practice short-answer and MCQ questions.
• Connect the concept with real-life applications.

Practice Questions:
1. What is the main idea of this topic?
2. List two important concepts from the notes.
3. Explain one real-life application.
4. Write a short exam answer on this topic.
5. Create one example based on this topic.

Final 3-Line Summary:
• The notes explain the main idea and related concepts.
• Students should focus on definitions, examples, and applications.
• Quick revision and quiz practice will improve understanding.""".format(
            overview=overview,
            points="\n".join(f"• {p}" for p in points[:5]) if points else "• Important concepts were not clearly detected.",
            terms="\n".join(f"• {t}: important term from the notes" for t in terms[:6]) if terms else "• Definition: meaning of the topic\n• Example: practical use of the topic",
        )

    def ai_generate_quiz(topic, source_text=""):
        topic = (topic or "General Learning").strip()
        context = safe_limit(source_text or "", 12000)
        prompt = f"""
Generate exactly 5 valid MCQ questions for students.

Topic: {topic}

Use the provided notes if available. Questions must be specific to the topic, not generic.

Return ONLY valid JSON in this format:
[
  {{
    "question": "...",
    "options": ["...", "...", "...", "..."],
    "answer_index": 0,
    "explanation": "..."
  }}
]

Rules:
• Each question must have 4 options.
• answer_index must be 0, 1, 2, or 3.
• Questions must be educationally meaningful.
• Do not generate generic questions like 'what is the first step'.

Notes:
{context}
"""
        ai_text = call_ai(prompt)
        if ai_text:
            try:
                match = re.search(r"\[.*\]", ai_text, re.S)
                data = json.loads(match.group(0) if match else ai_text)
                valid = []
                for item in data:
                    if isinstance(item, dict) and len(item.get("options", [])) == 4:
                        item["answer_index"] = int(item.get("answer_index", 0))
                        item["explanation"] = item.get("explanation", "This is the correct answer based on the topic.")
                        valid.append(item)
                if len(valid) >= 3:
                    return valid[:5]
            except Exception:
                pass
        return []

    def generate_questions(topic, source_text=""):
        ai_quiz = ai_generate_quiz(topic, source_text)
        if ai_quiz:
            return ai_quiz

        topic = (topic or "General Learning").strip()
        text = safe_limit(source_text or topic, 9000)
        lower = (topic + " " + text).lower()

        def q(question, options, answer, explanation):
            return {"question": question, "options": options, "answer_index": answer, "explanation": explanation}

        bank_topic = detect_topic(lower)
        if bank_topic == "dbms":
            return [
                q("Which key uniquely identifies each record in a table?", ["Primary Key", "Foreign Key", "Candidate Name", "HTML Tag"], 0, "A primary key uniquely identifies a row."),
                q("What does normalization mainly reduce?", ["Data redundancy", "Internet speed", "Screen brightness", "File size only"], 0, "Normalization reduces repeated data."),
                q("Which language is used to query relational databases?", ["SQL", "CSS", "HTML", "XML only"], 0, "SQL is used for database queries."),
                q("What is a foreign key used for?", ["Linking tables", "Changing font", "Running Python", "Compressing images"], 0, "Foreign keys create relationships between tables."),
                q("Which is an example of a DBMS?", ["SQLite", "MS Paint", "Keyboard", "Router cable"], 0, "SQLite is a database management system."),
            ]
        if bank_topic == "python":
            return [
                q("Which data structure stores ordered multiple values in Python?", ["List", "Router", "Table tag", "Voltage"], 0, "A list stores ordered values."),
                q("Which keyword defines a function in Python?", ["def", "func", "function", "method"], 0, "Python uses def."),
                q("Which symbol starts a comment in Python?", ["#", "//", "<!--", "@@"], 0, "# starts a Python comment."),
                q("Which function displays output?", ["print()", "show()", "displayHTML()", "echoCSS()"], 0, "print() outputs text."),
                q("What is Python commonly used for?", ["AI, web development, automation", "Only painting", "Only hardware wiring", "Only gaming controls"], 0, "Python has many uses."),
            ]
        if bank_topic == "ohm":
            return [
                q("What is Ohm's Law formula?", ["V = I × R", "F = m × a", "P = m × v", "A = πr²"], 0, "Ohm's Law is V=IR."),
                q("What does I represent in Ohm's Law?", ["Current", "Insulation", "Intensity only", "Input"], 0, "I represents current."),
                q("If I = 3A and R = 4Ω, what is voltage?", ["12V", "7V", "1.33V", "24V"], 0, "V=IR=3×4=12V."),
                q("Resistance is measured in which unit?", ["Ohm", "Volt", "Ampere", "Watt-hour"], 0, "Resistance is measured in ohms."),
                q("Ohm's Law is mainly used in", ["Electrical circuits", "Grammar", "History", "Cooking"], 0, "It is used in electrical circuits."),
            ]
        if bank_topic == "ai":
            return [
                q("What is Artificial Intelligence?", ["Machines performing human-like intelligent tasks", "Only storing files", "Only drawing images manually", "Only creating tables"], 0, "AI enables machines to do intelligent tasks."),
                q("Which is a subset of AI?", ["Machine Learning", "HTML", "Keyboard", "Database cable"], 0, "ML is a subset of AI."),
                q("Which AI field deals with human language?", ["Natural Language Processing", "Voltage Control", "File Compression", "Screen Design"], 0, "NLP handles human language."),
                q("Which is an AI application?", ["Chatbot", "Plain notebook", "Water bottle", "Wooden chair"], 0, "Chatbots use AI."),
                q("Why is AI useful in education?", ["It can support doubt-solving and personalization", "It removes all learning", "It prevents revision", "It only changes colors"], 0, "AI supports learning assistance."),
            ]
        return [
            q(f"Which statement best defines {topic}?", [f"{topic} is a study concept with definitions, examples, and applications", "It is only entertainment", "It is unrelated to learning", "It cannot be studied"], 0, "This checks the basic meaning."),
            q(f"What should students focus on while studying {topic}?", ["Definitions, concepts, and examples", "Only page numbers", "Only decoration", "Ignoring practice"], 0, "Core concepts and examples are important."),
            q(f"Which activity checks understanding of {topic}?", ["Attempting quizzes", "Closing notes", "Skipping examples", "Avoiding revision"], 0, "Quizzes test understanding."),
            q(f"Why are examples useful in {topic}?", ["They connect theory with practice", "They remove understanding", "They are only for design", "They are unrelated"], 0, "Examples improve clarity."),
            q(f"What is a good revision method for {topic}?", ["Summarize key points and practice questions", "Ignore notes", "Study only once", "Avoid self-test"], 0, "Summaries and practice help revision."),
        ]

    def suggest_questions_for_note(note):
        title = note["title"] if note else "this note"
        subject = note["subject"] if note else "the topic"
        return [
            f"What is the main concept of {title}?",
            f"Explain {subject} in simple words.",
            f"Give exam revision notes from {title}.",
            f"Create 5 meaningful quiz questions from {title}.",
        ]

    def generate_notes_text(topic):
        topic = topic.strip() or "General Topic"
        prompt = f"""
Create detailed teacher notes on {topic}.

Use this format:
# Title
## Introduction
## Key Concepts
## Important Definitions
## Examples
## Applications
## Quick Revision
## Summary

Write useful educational content, not generic filler.
"""
        ai_notes = call_ai(prompt)
        if ai_notes:
            return ai_notes
        return f"""# {topic}

## Introduction
{topic} is an important educational topic that helps learners understand concepts in a structured way.

## Key Concepts
• Definition and meaning of {topic}
• Important terms and principles
• Real-life examples
• Applications in academics and projects
• Revision and practice methods

## Detailed Explanation
To study {topic}, students should first understand the basic definition, then learn important subtopics and finally apply the concept through examples and practice questions.

## Example
A student can improve understanding of {topic} by reading notes, summarizing important points, asking doubts through the chatbot and attempting quizzes.

## Quick Revision
• Learn the definition.
• Revise keywords.
• Practice short questions.
• Attempt quizzes.

## Summary
{topic} should be learned through definitions, examples, applications and self-assessment."""


    # ---------- IQ200 useful feature overrides ----------
    def detect_topic(text: str) -> str:
        lower = (text or "").lower()
        if any(k in lower for k in ["aws", "amazon web services", "ec2", "s3", "lambda", "cloud computing"]):
            return "aws"
        if any(k in lower for k in ["ohm", "voltage", "current", "resistance", "circuit"]):
            return "ohm"
        if any(k in lower for k in ["dbms", "database", "sql", "normalization", "primary key", "foreign key"]):
            return "dbms"
        if any(k in lower for k in ["python", "list", "tuple", "dictionary", "programming"]):
            return "python"
        if any(k in lower for k in ["machine learning", "supervised", "unsupervised", "classification", "regression"]):
            return "ml"
        if any(k in lower for k in ["artificial intelligence", "deep learning", "neural", " ai "]):
            return "ai"
        if any(k in lower for k in ["data structure", "array", "stack", "queue", "linked list", "tree"]):
            return "ds"
        if any(k in lower for k in ["operating system", "process", "thread", "deadlock", "memory management"]):
            return "os"
        if any(k in lower for k in ["network", "tcp", "ip", "osi", "router", "dns"]):
            return "network"
        if any(k in lower for k in ["html", "css", "javascript", "react", "web development", "frontend", "backend"]):
            return "web"
        return "general"

    TOPIC_EXPLANATIONS = {
        "aws": """AWS (Amazon Web Services)

AWS is a cloud computing platform provided by Amazon. It allows users to rent computing power, storage, databases, networking, security tools, and AI services over the internet instead of buying physical servers.

Important services:
• EC2: virtual servers for running applications.
• S3: object storage for files, images, backups, and datasets.
• Lambda: serverless functions that run code without managing servers.
• RDS: managed relational databases.
• IAM: identity and access management for users and permissions.
• CloudFront: content delivery network for fast global delivery.

Why AWS is useful:
• It reduces infrastructure cost.
• It supports scalable applications.
• It is used for hosting websites, APIs, databases, AI apps, and enterprise systems.

Example:
A Flask website can be hosted on a cloud server, static files can be stored in S3, and user data can be stored in a managed database.""",
        "ai": """Artificial Intelligence (AI)

Artificial Intelligence is a branch of computer science that enables machines to perform tasks that normally require human intelligence, such as learning, reasoning, problem-solving, language understanding, and decision-making.

Key concepts:
• Machine Learning: systems learn patterns from data.
• Deep Learning: neural networks solve complex problems.
• Natural Language Processing: computers understand human language.
• Computer Vision: computers understand images and videos.
• Expert Systems: rule-based systems used for decision support.

Examples:
• Chatbots
• Recommendation systems
• Voice assistants
• Healthcare diagnosis
• Self-driving vehicles
• AI tutors in education

In education:
AI can help students by answering doubts, summarizing notes, generating quizzes, and supporting personalized learning.""",
        "ml": """Machine Learning

Machine Learning is a subset of Artificial Intelligence where computers learn patterns from data and improve performance without being explicitly programmed for every rule.

Types:
• Supervised Learning: learns from labeled data.
• Unsupervised Learning: finds patterns in unlabeled data.
• Reinforcement Learning: learns through rewards and penalties.

Examples:
• Spam email detection
• Product recommendations
• Image classification
• Student performance prediction

In education:
Machine learning can analyze student activity and suggest personalized learning resources.""",
        "dbms": """DBMS (Database Management System)

A DBMS is software used to store, organize, manage, secure, and retrieve data efficiently.

Key concepts:
• Table: stores data in rows and columns.
• Primary Key: uniquely identifies each row.
• Foreign Key: links one table to another.
• SQL: language used to query databases.
• Normalization: reduces duplicate data and improves consistency.
• Transaction: a safe unit of database work.

Examples:
SQLite, MySQL, PostgreSQL, Oracle, MongoDB.

Use in AccessLearn:
The database stores users, notes, quiz scores, summaries, meetings, and learning activity records.""",
        "python": """Python Programming

Python is a high-level programming language known for simple syntax and readability. It is widely used in web development, automation, data analysis, AI, and machine learning.

Key concepts:
• Variables store values.
• Lists and dictionaries store collections of data.
• Functions make code reusable.
• Loops repeat tasks.
• Modules and libraries extend functionality.

Example:
name = "AccessLearn"
print(name)

Why Python is useful:
It is beginner-friendly and has strong libraries for AI, Flask web development, data science, and automation.""",
        "ohm": """Ohm's Law

Ohm's Law explains the relationship between voltage, current, and resistance in an electrical circuit.

Formula:
V = I × R

Where:
• V = Voltage in volts
• I = Current in amperes
• R = Resistance in ohms

Example:
If current is 2 A and resistance is 5 Ω, then voltage = 2 × 5 = 10 V.

Use:
Ohm's Law is used to calculate voltage, current, or resistance in electrical and electronics circuits.""",
        "ds": """Data Structures

Data structures are ways of organizing data so it can be stored and used efficiently.

Important types:
• Array: stores elements in continuous memory.
• Stack: follows LIFO, last in first out.
• Queue: follows FIFO, first in first out.
• Linked List: stores data in connected nodes.
• Tree: represents hierarchical data.
• Graph: represents networks and connections.

Why important:
Data structures improve searching, sorting, memory usage, and algorithm performance.""",
        "os": """Operating System

An Operating System is system software that manages hardware and software resources and provides services to applications.

Important functions:
• Process management
• Memory management
• File management
• Device management
• Security
• User interface

Examples:
Windows, Linux, macOS, Android.

It acts as a bridge between the user, applications, and computer hardware.""",
        "network": """Computer Networking

Computer networking connects computers and devices so they can share data and resources.

Important concepts:
• IP Address: identifies a device on a network.
• Router: forwards data between networks.
• DNS: converts domain names to IP addresses.
• TCP/IP: communication protocol suite.
• OSI Model: explains network communication in layers.

Applications:
Internet, email, video calls, cloud services, and online learning platforms.""",
        "web": """Web Development

Web development is the process of building websites and web applications.

Main parts:
• Frontend: user interface using HTML, CSS, and JavaScript.
• Backend: server-side logic using Flask, Node.js, Django, etc.
• Database: stores application data.
• API: connects frontend and backend.

Example:
AccessLearn uses Flask as backend, SQLite as database, and HTML/CSS/JavaScript for the user interface."""
    }

    def topic_answer(topic: str, message: str = "") -> str:
        combined = (topic or "") + " " + (message or "")
        detected = detect_topic(combined)
        if detected in TOPIC_EXPLANATIONS:
            return TOPIC_EXPLANATIONS[detected]

        raw = (message or topic or "this topic").strip()
        cleaned = re.sub(r"^(what is|what is mean by|meaning of|explain|tell me about)\s+", "", raw, flags=re.I).strip(" ?.")
        label = cleaned.title() if cleaned else "This Topic"
        return f"""{label}

{label} refers to a concept, tool, subject, or process that should be understood through its definition, components, working, examples, and applications.

Student-friendly explanation:
• First understand what {label} means.
• Then identify its main parts or features.
• Study one real-world example.
• Learn where it is used.
• Practice questions to check understanding.

Important points:
• Definition gives the basic meaning.
• Key terms explain the structure of the topic.
• Examples make the topic easier to understand.
• Applications show practical value.
• Revision and quiz practice improve memory.

Note:
For deeper answers, add more specific detail in the question, such as "advantages of {label}", "types of {label}", or "{label} with examples"."""

    def local_explain(question: str, context: str = "") -> str:
        q = (question or "").strip()
        if not q:
            return "Please ask a clear question."

        # Special handling for notes-based answers
        if context and is_educational_content(context):
            relevance = note_relevance_score(q, context)
            if relevance > 0 or len(q.split()) <= 4:
                points = important_sentences(context, 7)
                terms = extract_keywords(context, 8)
                explanation = topic_answer(context, q)
                return """Notes-Based Answer

Question:
{question}

Answer:
{answer}

Important points from selected notes:
{points}

Important terms:
{terms}

How to revise:
• Read the answer once.
• Revise the important terms.
• Convert the points into short exam notes.
• Attempt a quiz based on this note.""".format(
                    question=q,
                    answer=explanation,
                    points="\n".join(f"• {p}" for p in points) if points else "• The selected note has limited extractable points.",
                    terms="\n".join(f"• {t}" for t in terms) if terms else "• Definition\n• Examples\n• Applications",
                )

        return topic_answer(q, q)

    def ai_answer(message, context=""):
        message = (message or "").strip()
        if not message:
            return "Please ask a clear question."

        usable_context = context if is_educational_content(context) else ""
        prompt = f"""
You are AccessLearn AI Assistant for students and teachers.

Answer the student's question accurately, clearly, and in detail.

Important rules:
1. Answer EVERY general and educational question directly.
2. Never return generic template text like "start with definition" unless the user asks for study methods.
3. If teacher notes are provided, analyze them and use them as context.
4. If selected notes are irrelevant, answer from general educational knowledge and mention briefly that the note was not sufficient.
5. Use headings, bullets, examples, and exam-friendly points.
6. For technical topics, include definition, key concepts, example, and applications.
7. Keep answer student-friendly.

Teacher Notes Context:
{safe_limit(usable_context, 12000)}

Student Question:
{message}
"""
        ai_text = call_ai(prompt)
        if ai_text:
            return ai_text
        return local_explain(message, usable_context)

    def summarize_text(text):
        clean_text = safe_limit(text, 22000)
        if not clean_text:
            return "Please enter notes to summarize."

        if not is_educational_content(clean_text):
            return """Content Quality Check

This content does not look like proper educational notes. It appears promotional, incomplete, or unrelated to study material.

Please upload/paste proper content such as:
• Teacher notes
• Textbook explanation
• Lecture content
• Topic-based PDF/DOC notes
• Definitions, concepts, examples, and applications

AccessLearn avoids creating misleading summaries from low-quality content."""

        prompt = f"""
You are AccessLearn Smart Notes Summarizer.

Read the complete notes carefully and create a useful student revision summary. Do NOT copy the first paragraph. Understand the content and rewrite it.

Required output format:

1. Topic Overview
- Explain the topic in 4-6 lines.

2. Key Concepts
- List 8-12 important concepts from the notes.

3. Important Definitions
- Define important terms in simple language.

4. Working / Process / Explanation
- Explain how the topic works or how the ideas are connected.

5. Applications / Examples
- Give practical examples and real-life uses.

6. Exam Revision Notes
- Write 8-12 bullet points useful for exams.

7. Common Mistakes / Confusions
- Mention common misunderstandings if relevant.

8. Practice Questions
- Generate 8 meaningful questions from the notes.

9. Final Short Summary
- Summarize the whole topic in 5 lines.

Notes:
{clean_text}
"""
        ai_summary = call_ai(prompt)
        if ai_summary:
            return ai_summary

        points = important_sentences(clean_text, 12)
        terms = extract_keywords(clean_text, 12)
        topic = terms[0] if terms else "This Topic"
        return """1. Topic Overview
{overview}

2. Key Concepts
{concepts}

3. Important Definitions
{definitions}

4. Working / Process / Explanation
• The topic should be understood by connecting its definition, features, examples, and applications.
• The selected notes contain multiple points that can be converted into short revision material.

5. Applications / Examples
• Academic learning
• Exam preparation
• Practical projects
• Real-world problem solving

6. Exam Revision Notes
{exam_points}

7. Common Mistakes / Confusions
• Do not memorize only keywords; understand meaning and examples.
• Do not skip applications because they help in long-answer questions.

8. Practice Questions
{practice}

9. Final Short Summary
• The notes explain {topic} and its related concepts.
• Students should focus on definitions, key points, examples, and applications.
• A short quiz after revision will help check understanding.""".format(
            overview=topic_answer(topic, topic).split("\n\n")[0],
            concepts="\n".join(f"• {t}" for t in terms[:10]) if terms else "• Key concepts were not clearly detected.",
            definitions="\n".join(f"• {t}: Important term related to this topic." for t in terms[:8]) if terms else "• Definition: basic meaning of the topic.",
            exam_points="\n".join(f"• {p}" for p in points[:10]) if points else "• Read the topic definition.\n• Revise examples.\n• Practice questions.",
            practice="\n".join(f"{i+1}. Explain {terms[i]} in simple words." for i in range(min(8, len(terms)))) if terms else "1. What is the main idea of this topic?\n2. List important concepts.\n3. Explain one example.\n4. Write applications.",
            topic=topic,
        )

    def build_quiz_bank(topic_text: str):
        detected = detect_topic(topic_text)
        def q(question, options, answer, explanation):
            return {"question": question, "options": options, "answer_index": answer, "explanation": explanation}

        banks = {
            "aws": [
                q("What does AWS stand for?", ["Amazon Web Services", "Advanced Web Storage", "Automated Wireless System", "Application Web Server"], 0, "AWS stands for Amazon Web Services."),
                q("Which AWS service is mainly used for object storage?", ["S3", "EC2", "IAM", "Route 53"], 0, "Amazon S3 stores objects such as files, images, and backups."),
                q("Which AWS service provides virtual servers?", ["EC2", "S3", "CloudFront", "DynamoDB"], 0, "EC2 provides scalable virtual machines."),
                q("What is IAM used for in AWS?", ["Managing users and permissions", "Editing images", "Writing CSS", "Making videos"], 0, "IAM controls identity and access permissions."),
                q("Which AWS service is serverless?", ["Lambda", "EC2", "EBS", "VPC"], 0, "AWS Lambda runs code without managing servers."),
                q("What is cloud computing?", ["Using computing resources over the internet", "Only using a local hard disk", "Drawing circuit diagrams", "Writing offline notes"], 0, "Cloud computing provides resources on demand via the internet."),
                q("Which service helps deliver content globally with low latency?", ["CloudFront", "IAM", "RDS", "SQS"], 0, "CloudFront is a content delivery network."),
                q("Which AWS service is commonly used for relational databases?", ["RDS", "S3", "Lambda", "CloudWatch"], 0, "RDS is a managed relational database service."),
                q("Why do companies use AWS?", ["Scalability and reduced infrastructure management", "Only for typing documents", "Only for graphic design", "Only for offline storage"], 0, "AWS helps scale apps and reduces hardware management."),
                q("Which AWS service is used to monitor logs and metrics?", ["CloudWatch", "EC2", "S3 Glacier", "VPC"], 0, "CloudWatch monitors metrics, logs, and alarms."),
            ],
            "ai": [
                q("What is Artificial Intelligence?", ["Machines performing human-like intelligent tasks", "Only storing files", "Only designing web pages", "A type of database key"], 0, "AI enables machines to learn, reason, and solve problems."),
                q("Which is a subset of AI?", ["Machine Learning", "HTML", "Router cable", "Primary key"], 0, "Machine Learning is a major subset of AI."),
                q("What does NLP help computers understand?", ["Human language", "Electric current only", "Hard disk speed", "CSS colors"], 0, "NLP focuses on human language processing."),
                q("Which AI field works with images and videos?", ["Computer Vision", "SQL", "Operating System", "Normalization"], 0, "Computer vision deals with visual data."),
                q("Which is an example of AI in daily life?", ["Voice assistant", "Plain notebook", "Wooden table", "Manual switch only"], 0, "Voice assistants use AI and NLP."),
                q("What is Deep Learning based on?", ["Neural networks", "HTML tags", "Primary keys", "Copper wires"], 0, "Deep learning uses artificial neural networks."),
                q("How can AI help students?", ["Doubt solving and personalized learning", "Only increasing screen brightness", "Only replacing books physically", "Only printing files"], 0, "AI can support explanations, summaries, and quizzes."),
                q("What is a training dataset?", ["Data used to teach a model", "A school timetable", "A power supply", "A CSS file"], 0, "Models learn patterns from training data."),
                q("What is model prediction?", ["Output generated by a trained model", "A router setting", "A document font", "A database table name"], 0, "Prediction is the model's answer/output."),
                q("Which risk is important in AI systems?", ["Bias and privacy issues", "Only paper color", "Only table border", "Only keyboard layout"], 0, "AI systems can involve bias and privacy concerns."),
            ],
            "dbms": [
                q("Which key uniquely identifies each record?", ["Primary Key", "Foreign Key", "CSS Selector", "IP Address"], 0, "A primary key uniquely identifies a row."),
                q("What does normalization mainly reduce?", ["Data redundancy", "Internet speed", "Screen brightness", "File extension"], 0, "Normalization reduces duplicate data."),
                q("Which language is used to query relational databases?", ["SQL", "HTML", "CSS", "XML only"], 0, "SQL is used to query relational databases."),
                q("What is a foreign key used for?", ["Linking tables", "Changing font", "Running Python loops", "Compressing video"], 0, "Foreign keys create relationships between tables."),
                q("Which is an example of a DBMS?", ["SQLite", "MS Paint", "Keyboard", "Router cable"], 0, "SQLite is a database management system."),
                q("What is a table?", ["Rows and columns storing data", "A cloud server", "A Python package only", "An image filter"], 0, "Tables store data in rows and columns."),
                q("What is a transaction?", ["A reliable unit of database work", "A CSS animation", "A web browser only", "A keyboard shortcut"], 0, "Transactions help maintain database consistency."),
                q("Which command is used to retrieve data?", ["SELECT", "STYLE", "PRINTSCREEN", "BOOT"], 0, "SELECT retrieves rows from a database."),
                q("Why is database security important?", ["To protect stored data", "To change monitor size", "To draw shapes", "To improve keyboard sound"], 0, "Security protects sensitive records."),
                q("What is data consistency?", ["Data remains accurate and valid", "Data becomes random", "Only images are stored", "Only fonts are changed"], 0, "Consistency ensures valid database state."),
            ],
            "python": [
                q("Which keyword defines a function in Python?", ["def", "func", "function", "method"], 0, "Python uses def."),
                q("Which data type stores ordered multiple values?", ["List", "Router", "HTML tag", "Voltage"], 0, "A list stores ordered values."),
                q("Which symbol starts a Python comment?", ["#", "//", "<!--", "@@"], 0, "# starts a Python comment."),
                q("Which function displays output?", ["print()", "show()", "displayHTML()", "echoCSS()"], 0, "print() outputs text."),
                q("What is a dictionary in Python?", ["Key-value collection", "Only a book", "A circuit formula", "A CSS property"], 0, "Dictionaries store key-value pairs."),
                q("Which loop is commonly used for iteration?", ["for", "repeatCSS", "selectSQL", "routeIP"], 0, "for loops iterate over sequences."),
                q("Why is Python popular for AI?", ["Strong libraries and simple syntax", "Only because of HTML", "It cannot use libraries", "It is only for painting"], 0, "Python has libraries like NumPy, pandas, and ML tools."),
                q("What is Flask?", ["Python web framework", "Database table", "Cloud storage", "Operating system"], 0, "Flask is a lightweight Python web framework."),
                q("What does pip do?", ["Installs Python packages", "Designs logos", "Controls voltage", "Manages routers only"], 0, "pip installs packages."),
                q("What is indentation used for?", ["Defining code blocks", "Changing screen brightness", "Creating database keys", "Sending emails only"], 0, "Python uses indentation to define blocks."),
            ],
            "ohm": [
                q("What is Ohm's Law formula?", ["V = I × R", "F = m × a", "P = m × v", "A = πr²"], 0, "Ohm's Law is V = I × R."),
                q("What does I represent?", ["Current", "Insulation", "Intensity only", "Input"], 0, "I represents current."),
                q("If I = 3A and R = 4Ω, voltage is", ["12V", "7V", "1.33V", "24V"], 0, "V = 3 × 4 = 12V."),
                q("Resistance is measured in", ["Ohm", "Volt", "Ampere", "Watt-hour"], 0, "Resistance is measured in ohms."),
                q("Ohm's Law is used in", ["Electrical circuits", "Grammar", "History", "Cooking"], 0, "It is used in electrical circuits."),
                q("What does V represent?", ["Voltage", "Volume only", "Velocity only", "Variable name only"], 0, "V represents voltage."),
                q("What happens to current if resistance increases and voltage is constant?", ["Current decreases", "Current increases", "Current becomes infinite", "No relation"], 0, "I = V/R, so current decreases."),
                q("Which instrument measures current?", ["Ammeter", "Voltmeter only", "Thermometer", "Barometer"], 0, "Ammeter measures current."),
                q("Which instrument measures voltage?", ["Voltmeter", "Ammeter only", "Scale", "Compass"], 0, "Voltmeter measures voltage."),
                q("Unit of current is", ["Ampere", "Ohm", "Volt", "Joule"], 0, "Current is measured in amperes."),
            ],
        }
        if detected in banks:
            return banks[detected]
        return []

    def ai_generate_quiz(topic, source_text="", count: int = 10):
        topic = (topic or "General Learning").strip()
        context = safe_limit(source_text or "", 14000)
        prompt = f"""
Generate exactly {count} valid MCQ questions for students.

Topic: {topic}

Use the provided notes if available. Questions must be specific to the topic and conceptually valid.

Return ONLY valid JSON in this format:
[
  {{
    "question": "...",
    "options": ["...", "...", "...", "..."],
    "answer_index": 0,
    "explanation": "..."
  }}
]

Rules:
• Each question must have 4 options.
• answer_index must be 0, 1, 2, or 3.
• Questions must not be generic.
• Questions should test definitions, concepts, applications, examples, differences, and reasoning.
• Avoid repeating the topic name in every question.
• Explanations should be short and useful.

Notes:
{context}
"""
        ai_text = call_ai(prompt)
        if ai_text:
            try:
                match = re.search(r"\[.*\]", ai_text, re.S)
                data = json.loads(match.group(0) if match else ai_text)
                valid = []
                seen = set()
                for item in data:
                    if isinstance(item, dict) and len(item.get("options", [])) == 4:
                        question = str(item.get("question", "")).strip()
                        if not question or question.lower() in seen:
                            continue
                        seen.add(question.lower())
                        item["answer_index"] = max(0, min(3, int(item.get("answer_index", 0))))
                        item["explanation"] = item.get("explanation", "This is the correct answer based on the topic.")
                        valid.append(item)
                if len(valid) >= 5:
                    return valid[:count]
            except Exception:
                pass
        return []

    def generate_questions(topic, source_text=""):
        topic = (topic or "").strip()
        context = safe_limit(source_text or "", 12000)
        ai_quiz = ai_generate_quiz(topic or "Selected Notes", context, 10)
        if ai_quiz:
            return ai_quiz

        combined = (topic + " " + context).strip()
        bank = build_quiz_bank(combined)
        if bank:
            return bank[:10]

        # Last fallback: create questions from extracted terms, not generic "first step" questions.
        terms = extract_keywords(combined, 12)
        if not terms:
            terms = [topic or "the topic", "definition", "application", "example", "concept"]
        qs = []
        def q(question, options, answer, explanation):
            return {"question": question, "options": options, "answer_index": answer, "explanation": explanation}
        main = terms[0]
        qs.append(q(f"What is the main meaning of {main}?", [f"Core concept related to {main}", "Random unrelated term", "Only a file name", "Only a color"], 0, f"{main} is the central idea of the selected topic."))
        qs.append(q(f"Why is {main} important in learning?", ["It helps understand the subject clearly", "It has no use", "It only changes font", "It only stores images"], 0, "Understanding importance helps connect theory with use."))
        for term in terms[1:9]:
            qs.append(q(f"Which statement best describes {term}?", [f"{term} is an important concept in this topic", f"{term} is unrelated to the topic", f"{term} is only a design color", f"{term} is only a keyboard key"], 0, f"{term} was detected as an important term in the content."))
        return qs[:10]

    def suggest_questions_for_note(note):
        if not note:
            return [
                "What is Artificial Intelligence?",
                "Explain AWS in simple words.",
                "Difference between RAM and ROM.",
                "Generate exam notes on DBMS.",
                "Create 10 quiz questions on Python."
            ]
        title = note["title"] or "this note"
        subject = note["subject"] or "the topic"
        text, _ = note_context(note["id"])
        terms = extract_keywords(text + " " + title + " " + subject, 5)
        base = terms[:3] if terms else [subject, title]
        questions = [
            f"Explain {subject} in simple words.",
            f"What are the key concepts in {title}?",
            f"Give exam revision notes from {title}.",
            f"Create 10 meaningful quiz questions from {title}.",
        ]
        for term in base:
            questions.append(f"What is {term} and why is it important?")
        return questions[:7]


    @app.route("/")
    def index():
        if session.get("user_id"):
            return redirect(url_for("dashboard"))
        return render_template("index.html")

    @app.route("/register", methods=["GET", "POST"])
    def register():
        if request.method == "POST":
            name = request.form.get("name", "").strip()
            email = request.form.get("email", "").strip().lower()
            password = request.form.get("password", "").strip()
            role = request.form.get("role", "student")
            if not name or not email or not password or role not in ("student", "teacher"):
                flash("Please fill all fields correctly.", "danger")
                return render_template("register.html")
            db = get_db()
            try:
                cursor = db.execute("INSERT INTO users (name, email, password_hash, role) VALUES (?, ?, ?, ?)", (name, email, generate_password_hash(password), role))
                db.commit()
            except Exception:
                flash("Email already exists or registration failed.", "danger")
                return render_template("register.html")
            session["user_id"] = cursor.lastrowid
            session["role"] = role
            flash("Registration successful.", "success")
            return redirect(url_for("dashboard"))
        return render_template("register.html")

    @app.route("/login", methods=["GET", "POST"])
    def login():
        if request.method == "POST":
            email = request.form.get("email", "").strip().lower()
            password = request.form.get("password", "").strip()
            user = get_db().execute("SELECT * FROM users WHERE email = ?", (email,)).fetchone()
            if user and check_password_hash(user["password_hash"], password):
                session["user_id"] = user["id"]
                session["role"] = user["role"]
                flash("Login successful.", "success")
                return redirect(url_for("dashboard"))
            flash("Invalid email or password.", "danger")
        return render_template("login.html")

    @app.route("/logout")
    def logout():
        session.clear()
        flash("Logged out successfully.", "success")
        return redirect(url_for("index"))

    @app.route("/dashboard")
    @login_required
    def dashboard():
        user = current_user()
        db = get_db()
        total_notes = db.execute("SELECT COUNT(*) AS c FROM notes").fetchone()["c"]
        recent_notes = db.execute("SELECT n.*, u.name AS teacher_name FROM notes n LEFT JOIN users u ON n.teacher_id = u.id ORDER BY n.created_at DESC LIMIT 5").fetchall()
        meetings = db.execute("SELECT m.*, u.name AS teacher_name FROM meetings m LEFT JOIN users u ON m.teacher_id=u.id ORDER BY m.meeting_date DESC, m.meeting_time DESC LIMIT 5").fetchall()
        if user["role"] == "teacher":
            teacher_notes = db.execute("SELECT COUNT(*) AS c FROM notes WHERE teacher_id = ?", (user["id"],)).fetchone()["c"]
            total_students = db.execute("SELECT COUNT(*) AS c FROM users WHERE role='student'").fetchone()["c"]
            total_study_seconds = db.execute("SELECT COALESCE(SUM(time_spent),0) AS s FROM study_logs sl JOIN notes n ON sl.note_id=n.id WHERE n.teacher_id=?", (user["id"],)).fetchone()["s"]
            total_published_quizzes = db.execute("SELECT COUNT(*) AS c FROM published_quizzes WHERE teacher_id=?", (user["id"],)).fetchone()["c"]
            return render_template("teacher_dashboard.html", total_notes=total_notes, teacher_notes=teacher_notes, recent_notes=recent_notes, total_students=total_students, total_study_minutes=round(total_study_seconds/60, 1), meetings=meetings, total_published_quizzes=total_published_quizzes)
        total_quizzes = db.execute("SELECT COUNT(*) AS c FROM quiz_results WHERE student_id = ?", (user["id"],)).fetchone()["c"]
        avg_score_row = db.execute("SELECT AVG(CAST(score AS FLOAT)/total_questions) AS avg FROM quiz_results WHERE student_id = ?", (user["id"],)).fetchone()
        avg_score = round((avg_score_row["avg"] or 0) * 100, 1)
        recent_results = db.execute("SELECT * FROM quiz_results WHERE student_id = ? ORDER BY created_at DESC LIMIT 5", (user["id"],)).fetchall()
        published_quizzes = db.execute("SELECT pq.*, u.name AS teacher_name FROM published_quizzes pq LEFT JOIN users u ON pq.teacher_id=u.id ORDER BY pq.created_at DESC LIMIT 5").fetchall()
        return render_template("student_dashboard.html", total_notes=total_notes, total_quizzes=total_quizzes, avg_score=avg_score, recent_notes=recent_notes, recent_results=recent_results, meetings=meetings, published_quizzes=published_quizzes)

    @app.route("/notes")
    @login_required
    def notes():
        rows = get_db().execute("SELECT n.*, u.name AS teacher_name FROM notes n LEFT JOIN users u ON n.teacher_id = u.id ORDER BY n.created_at DESC").fetchall()
        return render_template("notes.html", notes=rows)

    @app.route("/notes/<int:note_id>")
    @login_required
    def note_detail(note_id):
        note = get_db().execute("SELECT n.*, u.name AS teacher_name FROM notes n LEFT JOIN users u ON n.teacher_id=u.id WHERE n.id=?", (note_id,)).fetchone()
        if not note:
            flash("Note not found.", "danger")
            return redirect(url_for("notes"))
        suggestions = suggest_questions_for_note(note)
        return render_template("note_detail.html", note=note, suggestions=suggestions)

    @app.route("/teacher/upload", methods=["GET", "POST"])
    @login_required
    @role_required("teacher")
    def upload_notes():
        if request.method == "POST":
            title = request.form.get("title", "").strip()
            subject = request.form.get("subject", "").strip()
            description = request.form.get("description", "").strip()
            content = request.form.get("content", "").strip()
            resource_link = request.form.get("resource_link", "").strip()
            video_link = request.form.get("video_link", "").strip()
            uploaded = request.files.get("resource_file")
            file_filename = file_original_name = file_type = None
            if uploaded and uploaded.filename:
                if not allowed_file(uploaded.filename):
                    flash("Unsupported file type.", "danger")
                    return render_template("upload_notes.html")
                file_original_name = uploaded.filename
                ext = uploaded.filename.rsplit('.', 1)[1].lower()
                file_type = ext
                file_filename = f"{datetime.now().strftime('%Y%m%d%H%M%S')}_{secure_filename(uploaded.filename)}"
                uploaded.save(UPLOAD_DIR / file_filename)
            if not title or not subject or not (content or file_filename or resource_link or video_link):
                flash("Title, subject, and at least one resource are required.", "danger")
                return render_template("upload_notes.html")
            db = get_db()
            db.execute("INSERT INTO notes (teacher_id, title, subject, description, content, file_filename, file_original_name, file_type, resource_link, video_link) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)", (session["user_id"], title, subject, description, content, file_filename, file_original_name, file_type, resource_link, video_link))
            db.commit()
            flash("Learning resource uploaded successfully.", "success")
            return redirect(url_for("notes"))
        return render_template("upload_notes.html")

    @app.route("/uploads/<path:filename>")
    @login_required
    def uploaded_file(filename):
        return send_from_directory(app.config["UPLOAD_FOLDER"], filename)

    @app.route("/teacher/delete/<int:note_id>")
    @login_required
    @role_required("teacher")
    def delete_note(note_id):
        db = get_db()
        db.execute("DELETE FROM notes WHERE id = ? AND teacher_id = ?", (note_id, session["user_id"]))
        db.commit()
        flash("Note deleted successfully.", "success")
        return redirect(url_for("notes"))

    @app.route("/teacher/edit/<int:note_id>", methods=["GET", "POST"])
    @login_required
    @role_required("teacher")
    def edit_note(note_id):
        db = get_db()
        note = db.execute("SELECT * FROM notes WHERE id = ? AND teacher_id = ?", (note_id, session["user_id"])).fetchone()
        if not note:
            flash("Note not found.", "danger")
            return redirect(url_for("notes"))
        if request.method == "POST":
            title = request.form.get("title", "").strip()
            subject = request.form.get("subject", "").strip()
            description = request.form.get("description", "").strip()
            content = request.form.get("content", "").strip()
            resource_link = request.form.get("resource_link", "").strip()
            video_link = request.form.get("video_link", "").strip()
            db.execute("UPDATE notes SET title=?, subject=?, description=?, content=?, resource_link=?, video_link=? WHERE id=? AND teacher_id=?", (title, subject, description, content, resource_link, video_link, note_id, session["user_id"]))
            db.commit()
            flash("Note updated successfully.", "success")
            return redirect(url_for("notes"))
        return render_template("edit_note.html", note=note)

    @app.route("/teacher/reports")
    @login_required
    @role_required("teacher")
    def teacher_reports():
        db = get_db()
        study_logs = db.execute("""
            SELECT sl.*, u.name AS student_name, n.title AS note_title, n.subject AS note_subject
            FROM study_logs sl
            JOIN users u ON sl.student_id = u.id
            JOIN notes n ON sl.note_id = n.id
            WHERE n.teacher_id = ?
            ORDER BY sl.viewed_at DESC
        """, (session["user_id"],)).fetchall()
        quiz_attempts = db.execute("""
            SELECT qa.*, u.name AS student_name, pq.title AS quiz_title, pq.topic
            FROM quiz_attempts qa
            JOIN users u ON qa.student_id=u.id
            JOIN published_quizzes pq ON qa.quiz_id=pq.id
            WHERE pq.teacher_id=?
            ORDER BY qa.attempted_at DESC
        """, (session["user_id"],)).fetchall()
        summary = db.execute("""
            SELECT
              COUNT(DISTINCT sl.student_id) AS active_students,
              COUNT(DISTINCT sl.note_id) AS viewed_notes,
              COALESCE(SUM(sl.time_spent),0) AS total_seconds
            FROM study_logs sl
            JOIN notes n ON sl.note_id=n.id
            WHERE n.teacher_id=?
        """, (session["user_id"],)).fetchone()
        quiz_summary = db.execute("""
            SELECT COUNT(*) AS attempts,
                   COALESCE(AVG(CAST(score AS FLOAT)/total_questions),0) AS avg_score
            FROM quiz_attempts qa
            JOIN published_quizzes pq ON qa.quiz_id=pq.id
            WHERE pq.teacher_id=?
        """, (session["user_id"],)).fetchone()
        return render_template(
            "teacher_reports.html",
            logs=study_logs,
            attempts=quiz_attempts,
            active_students=summary["active_students"] or 0,
            viewed_notes=summary["viewed_notes"] or 0,
            total_minutes=round((summary["total_seconds"] or 0)/60, 1),
            quiz_attempt_count=quiz_summary["attempts"] or 0,
            avg_quiz_score=round((quiz_summary["avg_score"] or 0)*100, 1),
        )

    @app.route("/meetings", methods=["GET", "POST"])
    @login_required
    def meetings():
        db = get_db()
        user = current_user()
        if request.method == "POST":
            if user["role"] != "teacher":
                flash("Only teachers can create meetings.", "danger")
                return redirect(url_for("meetings"))
            title = request.form.get("title", "").strip()
            meeting_date = request.form.get("meeting_date", "").strip()
            meeting_time = request.form.get("meeting_time", "").strip()
            meeting_link = request.form.get("meeting_link", "").strip()
            description = request.form.get("description", "").strip()
            if not title or not meeting_link:
                flash("Meeting title and link are required.", "danger")
                return redirect(url_for("meetings"))
            db.execute("INSERT INTO meetings (teacher_id, title, meeting_date, meeting_time, meeting_link, description) VALUES (?, ?, ?, ?, ?, ?)", (user["id"], title, meeting_date, meeting_time, meeting_link, description))
            db.commit()
            flash("Meeting created successfully.", "success")
            return redirect(url_for("meetings"))
        rows = db.execute("SELECT m.*, u.name AS teacher_name FROM meetings m LEFT JOIN users u ON m.teacher_id=u.id ORDER BY m.meeting_date DESC, m.meeting_time DESC").fetchall()
        return render_template("meetings.html", meetings=rows)

    # ---------- teacher quizzes ----------
    @app.route("/teacher/quizzes", methods=["GET", "POST"])
    @login_required
    @role_required("teacher")
    def teacher_quizzes():
        db = get_db()
        notes_rows = db.execute("SELECT * FROM notes WHERE teacher_id=? ORDER BY created_at DESC", (session["user_id"],)).fetchall()
        if request.method == "POST":
            title = request.form.get("title", "").strip()
            topic = request.form.get("topic", "").strip()
            note_id = int(request.form.get("note_id", 0) or 0)
            source_text = ""
            if note_id:
                source_text, note_title = note_context(note_id)
                if not topic:
                    topic = note_title
            if not title:
                title = f"{topic or 'Learning'} Quiz"
            questions = generate_questions(topic or title, source_text)
            db.execute("INSERT INTO published_quizzes (teacher_id, title, topic, note_id, questions_json) VALUES (?, ?, ?, ?, ?)", (session["user_id"], title, topic or title, note_id or None, json.dumps(questions)))
            db.commit()
            flash("Quiz generated and published successfully.", "success")
            return redirect(url_for("teacher_quizzes"))
        quizzes = db.execute("SELECT pq.*, n.title AS note_title FROM published_quizzes pq LEFT JOIN notes n ON pq.note_id=n.id WHERE pq.teacher_id=? ORDER BY pq.created_at DESC", (session["user_id"],)).fetchall()
        return render_template("teacher_quizzes.html", notes=notes_rows, quizzes=quizzes)

    @app.route("/quizzes")
    @login_required
    def published_quizzes():
        rows = get_db().execute("SELECT pq.*, u.name AS teacher_name FROM published_quizzes pq LEFT JOIN users u ON pq.teacher_id=u.id ORDER BY pq.created_at DESC").fetchall()
        return render_template("published_quizzes.html", quizzes=rows)

    @app.route("/quiz/published/<int:quiz_id>", methods=["GET", "POST"])
    @login_required
    def attempt_published_quiz(quiz_id):
        db = get_db()
        quiz = db.execute("SELECT pq.*, u.name AS teacher_name FROM published_quizzes pq LEFT JOIN users u ON pq.teacher_id=u.id WHERE pq.id=?", (quiz_id,)).fetchone()
        if not quiz:
            flash("Quiz not found.", "danger")
            return redirect(url_for("published_quizzes"))
        questions = json.loads(quiz["questions_json"])
        if request.method == "POST":
            score = 0
            answers = []
            for i, q in enumerate(questions):
                selected = request.form.get(f"q{i}")
                selected_index = int(selected) if selected is not None and selected.isdigit() else -1
                is_correct = selected_index == int(q.get("answer_index", 0))
                if is_correct:
                    score += 1
                answers.append({"selected": selected_index, "correct": q.get("answer_index", 0), "is_correct": is_correct})
            db.execute("INSERT INTO quiz_attempts (quiz_id, student_id, score, total_questions, answers_json) VALUES (?, ?, ?, ?, ?)", (quiz_id, session["user_id"], score, len(questions), json.dumps(answers)))
            db.execute("INSERT INTO quiz_results (student_id, username, topic, score, total_questions) VALUES (?, ?, ?, ?, ?)", (session["user_id"], current_user()["name"], quiz["topic"], score, len(questions)))
            db.commit()
            return render_template("attempt_quiz.html", quiz=quiz, questions=questions, score=score, submitted=True)
        return render_template("attempt_quiz.html", quiz=quiz, questions=questions, submitted=False)

    @app.route("/teacher/quiz-reports")
    @login_required
    @role_required("teacher")
    def teacher_quiz_reports():
        rows = get_db().execute("""
            SELECT qa.*, u.name AS student_name, pq.title AS quiz_title, pq.topic
            FROM quiz_attempts qa
            JOIN users u ON qa.student_id=u.id
            JOIN published_quizzes pq ON qa.quiz_id=pq.id
            WHERE pq.teacher_id=?
            ORDER BY qa.attempted_at DESC
        """, (session["user_id"],)).fetchall()
        return render_template("teacher_quiz_reports.html", attempts=rows)

    # ---------- APIs ----------
    @app.post("/api/study-log")
    @login_required
    def api_study_log():
        data = request.get_json(silent=True) or request.form or {}
        note_id = int(data.get("note_id", 0) or 0)
        time_spent = int(data.get("time_spent", 0) or 0)
        if note_id > 0 and time_spent > 0:
            get_db().execute("INSERT INTO study_logs (student_id, note_id, time_spent) VALUES (?, ?, ?)", (session["user_id"], note_id, time_spent))
            get_db().commit()
        return jsonify({"success": True})

    @app.post("/api/generate_notes")
    @login_required
    def generate_notes():
        data = request.get_json(force=True) or {}
        topic = data.get("topic", "").strip()
        if not topic:
            return jsonify({"notes": "Please enter a topic."})
        return jsonify({"notes": generate_notes_text(topic)})

    @app.get("/api/note-suggestions/<int:note_id>")
    @login_required
    def api_note_suggestions(note_id):
        note = get_db().execute("SELECT * FROM notes WHERE id=?", (note_id,)).fetchone()
        return jsonify({"questions": suggest_questions_for_note(note)})

    @app.route("/chatbot")
    @login_required
    def chatbot_page():
        notes_rows = get_db().execute("SELECT id, title, subject FROM notes ORDER BY created_at DESC").fetchall()
        return render_template("chatbot.html", notes=notes_rows)

    @app.route("/summarizer")
    @login_required
    def summarizer_page():
        notes_rows = get_db().execute("SELECT id, title, subject FROM notes ORDER BY created_at DESC").fetchall()
        return render_template("summarizer.html", notes=notes_rows)

    @app.route("/quiz")
    @login_required
    def quiz_page():
        notes_rows = get_db().execute("SELECT id, title, subject FROM notes ORDER BY created_at DESC").fetchall()
        return render_template("quiz.html", notes=notes_rows)

    @app.post("/api/chat")
    @login_required
    def api_chat():
        data = request.get_json(force=True) or {}
        message = data.get("message", "")
        note_id = int(data.get("note_id", 0) or 0)
        context, _ = note_context(note_id)
        reply = ai_answer(message, context)
        db = get_db()
        db.execute("INSERT INTO chat_logs (user_id, note_id, message, response) VALUES (?, ?, ?, ?)", (session["user_id"], note_id or None, message, reply))
        db.commit()
        return jsonify({"reply": reply})

    @app.post("/api/summarize")
    @login_required
    def api_summarize():
        source_type = "text"
        source_title = "Pasted Text"
        text = ""
        if request.content_type and "multipart/form-data" in request.content_type:
            source_type = request.form.get("source_type", "text")
            text = request.form.get("text", "")
            url = request.form.get("url", "").strip()
            note_id = int(request.form.get("note_id", 0) or 0)
            uploaded = request.files.get("file")
            if note_id:
                text, source_title = note_context(note_id)
                source_type = "uploaded_note"
            elif uploaded and uploaded.filename:
                tmp_name = f"summary_{datetime.now().strftime('%Y%m%d%H%M%S')}_{secure_filename(uploaded.filename)}"
                tmp_path = UPLOAD_DIR / tmp_name
                uploaded.save(tmp_path)
                text = extract_file_text(tmp_name)
                source_type = "file"
                source_title = uploaded.filename
                try:
                    tmp_path.unlink()
                except Exception:
                    pass
            elif url:
                text = fetch_link_text(url)
                source_type = "link"
                source_title = url
        else:
            data = request.get_json(force=True) or {}
            text = data.get("text", "")
            note_id = int(data.get("note_id", 0) or 0)
            if note_id:
                text, source_title = note_context(note_id)
                source_type = "uploaded_note"
        if not text:
            summary = "Could not read content. Please paste text, select a note, upload a supported PDF/DOCX/PPTX/TXT file, or use a readable webpage link."
        else:
            summary = summarize_text(text)
        db = get_db()
        db.execute("INSERT INTO summaries (user_id, source_type, source_title, original_text, summary) VALUES (?, ?, ?, ?, ?)", (session["user_id"], source_type, source_title, safe_limit(text, 2000), summary))
        db.commit()
        return jsonify({"summary": summary})

    @app.post("/api/generate_quiz")
    @login_required
    def api_generate_quiz():
        data = request.get_json(force=True) or {}
        topic = data.get("topic", "General Learning")
        note_id = int(data.get("note_id", 0) or 0)
        context, note_title = note_context(note_id)
        if note_id and not topic:
            topic = note_title
        return jsonify({"questions": generate_questions(topic, context)})

    @app.post("/api/submit_quiz")
    @login_required
    def api_submit_quiz():
        data = request.get_json(force=True) or {}
        topic = data.get("topic", "General Learning")
        score = int(data.get("score", 0))
        total_questions = int(data.get("total_questions", 0))
        user = current_user()
        db = get_db()
        db.execute("INSERT INTO quiz_results (student_id, username, topic, score, total_questions) VALUES (?, ?, ?, ?, ?)", (user["id"], user["name"], topic, score, total_questions))
        db.commit()
        return jsonify({"success": True})

    return app


app = create_app()

if __name__ == "__main__":
    port = int(os.getenv("PORT", "5000"))
    app.run(host="0.0.0.0", port=port, debug=True)
